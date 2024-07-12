import argparse
import os
import logging
from dotenv import load_dotenv
from PIL import Image
import pytesseract
import fitz
from docx import Document
from bs4 import BeautifulSoup
import pandas as pd
from pptx import Presentation
import re 

# Load environment variables
load_dotenv()

# Access environment variables
graphrag_api_key = os.getenv('OPENAI_API_KEY')

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def file_exists(file_path):
    if os.path.isfile(file_path):
        print(f"File '{file_path}' already exists.")
        return True
    else:
        print(f"File '{file_path}' does not exist. Creating it...")
        return False

def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as pdf_document:
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            text += page.get_text("text")
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'lxml')  # Use 'lxml' as the parser
    
    for element in soup.find_all(['p', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'div']):
        element.append('\n')
    
    # Get the raw text
    text = soup.get_text()
    
    # Remove leading/trailing whitespace and condense multiple newlines
    text = text.strip()
    text = re.sub(r'\n\s*\n+', '\n\n', text)
    
    return text

def extract_text_from_xlsx(file_path):
    df = pd.read_excel(file_path, sheet_name=None)
    text = ""
    for sheet_name, sheet_df in df.items():
        text += f"Sheet: {sheet_name}\n"
        text += sheet_df.to_string(index=False)
        text += "\n"
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_document(file_path):
    if file_path.endswith('.pdf'):
        # loader = PyPDFLoader(file_path)
        # documents = loader.load()
        # return "\n".join([doc.page_content for doc in documents])
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    elif file_path.endswith('.html'):
        return extract_text_from_html(file_path)
    elif file_path.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.heif')):
        # Use pytesseract to extract text from image files
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.xlsx'):
        return extract_text_from_xlsx(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f'Unsupported file format: {file_path}')

def preprocess_files(old_file_dir, new_file_dir):
    for root, _, files in os.walk(old_file_dir):
        for file in files:
            old_file_path = os.path.join(root, file)
            new_file_path = os.path.join(new_file_dir, f"{os.path.splitext(file)[0]}.txt")
            print(f"old_file_path: {old_file_path}")
            print(f"new_file_path: {new_file_path}")
            try:
                if file_exists(new_file_path):
                    continue
                else:
                    text = load_document(old_file_path)
                    with open(new_file_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                    logger.info(f"Processed and saved {file} to {new_file_path}")
            except Exception as e:
                logger.error(f"Error processing file {file}: {e}")

def run_indexing(subfolder):
    new_file_dir = os.path.join(subfolder, 'input') # final directory to store postprocessed files
    old_file_dir = os.path.join(subfolder, 'temp_input') # temp directory to store preprocessed files
    
    # Preprocess files
    preprocess_files(old_file_dir, new_file_dir)
    
    # Run the GraphRAG indexing pipeline
    os.system(f"python -m graphrag.index --root {subfolder}")

    # Example logging statement
    logger.info("Indexing completed successfully")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Run GraphRAG indexing pipeline.')
    parser.add_argument('subfolder', type=str, help='Subfolder to index (e.g., general, level_a, admin)')
    
    args = parser.parse_args()
    run_indexing(args.subfolder)


# Run the script examples:
# python script.py general
# python script.py level_q
# python script.py admin
